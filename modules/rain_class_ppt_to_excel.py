import re
import os
import pandas as pd
from pptx import Presentation

exerciseNum = 0


def read_shape_text(shape):
    '''获取形状中的文本'''
    text_runs = []
    if shape.shape_type == 6:  # type 6 是组合图形
        for s in shape.shapes:
            text_runs.extend(read_shape_text(s))  # 递归地读取组合图形内的每个形状
    elif shape.has_text_frame:  # 如果当前图像包含文本
        for paragraph in shape.text_frame.paragraphs:  # 遍历每一段
            para_text = ''
            for run in paragraph.runs:  # 遍历每一段中的每一行
                if run.text not in ['设置', '提交']:  # 如果当前行不是设置或提交
                    para_text += run.text  # 将当前行的文本添加到para_text中
            if para_text:  # 如果para_text不为空
                text_runs.append(para_text)  # 将para_text添加到text_runs中
    return text_runs


def read_pptx_text(file, filename):
    if 'ppt' not in file.split('.')[2].lower():
        print('\"' + file.split('PPT/')[1] + '\" ' + '不是PowerPoint演示文稿，跳过\n')
        return ''

    if 'pptx' not in file.split('.')[2].lower():
        print('\"' + file.split('PPT/')[1] + '\" ' + '格式已过时，跳过，请使用PowerPoint或WPS将其另存为.pptx格式\n')
        return ''

    prs = Presentation(file)
    text_runs = ['\n# ' + filename]  # 初始化一个列表，用于存储每个形状内的文本
    global exerciseNum
    exercise = 0
    exercise_page = []
    for slide in prs.slides:  # 遍历每一页
        for shape in slide.shapes:  # 遍历每一页的每一个shape
            if shape.shape_type == 6:  # type 6 是组合图形
                if shape.name == '组合 7':
                    exercise += 1
                    exercise_page.append(slide.slide_id)

    if exercise == 0:
        print('\"' + file.split('PPT/')[1] + '\" ' + '不是或没有包含标准格式的雨课堂习题，跳过\n')
        return ''
    else:
        exerciseNum += exercise
        print('该PPT包含 ' + str(exercise) + ' 道雨课堂习题')
        print('习题所在页的编号为：' + str(exercise_page))

    exercise_page.reverse()

    temp = 1
    for slide in prs.slides:  # 遍历每一页
        if slide.slide_id == exercise_page[len(exercise_page) - 1]:
            # print('正在获取第' + str(temp) + '道习题，共' + str(exercise) + '道')
            temp += 1
            exercise_page.pop()
            for shape in slide.shapes:  # 遍历每一页的每一个shape
                if shape.shape_type == 6:  # type 6 是组合图形
                    text_runs.extend(read_shape_text(shape))  # 递归地读取组合图形内的每个形状
                    break
                if not shape.has_text_frame:  # 如果当前图像不包含文本，跳过
                    continue
                if shape.fill.type == 1 and shape.fill.fore_color.rgb == (0, 255, 0):  # 如果当前形状有填充色，且填充色是绿色，为该项添加“答案“标签
                    text_runs.append('答案')
                text_runs.extend(read_shape_text(shape))  # 其他情况，直接读取图形内的文本
    print('\"' + file.split('PPT/')[1] + '\" ' + '读取完成\n')
    return '\n'.join(text_runs)


def format_text(text):
    # 将文本分成多行
    lines = text.replace(u'\xa0', '').split('\n')
    # 初始化变量
    question = ''
    options = []
    answer = ''
    i = 0
    k = 0
    result = ''

    # 遍历每一行
    for line in lines:
        if line.startswith('#'):
            # 如果存在现有问题,则将其格式化并添加到结果中
            if question:
                result += question + '\n'
                for option in options:
                    result += option + '\n'
                result += '答案：' + answer + '\n\n'
                # 重置变量
                question = ''
                options.clear()
                answer = ''

            result += line + '\n'
            i += 1
            k = 0
            continue
        # 检查该行是否为问题
        if re.search('\d+\.', line):
            k += 1
            # 如果存在现有问题,则将其格式化并添加到结果中
            if question:
                result += question + '\n'
                for option in options:
                    result += option + '\n'
                result += '答案：' + answer + '\n\n'
                # 重置变量
                question = ''
                options.clear()
                answer = ''

            j = i + 5
            while True:
                if '单选题' in lines[j] or '多选题' in lines[j]:
                    # 将当前行设置为问题
                    question = '(' + lines[j].split('题')[0] + ') ' + str(k) + '.' + line.split('.')[-1]
                    break
                j += 1

        # 检查该行是否为选项
        if re.search('[A-Z]', line):
            # 将选项添加到选项列表中
            if '答案' in lines[i - 1]:
                options.append(line + '.' + lines[i - 2])
            else:
                options.append(line + '.' + lines[i - 1])

        # 检查该行是否为答案
        if '答案' in line:
            # 设置答案
            answer += lines[i + 1]
        i += 1

    # 格式化并将最后一个问题添加到结果中
    if question:
        result += question + '\n'
        for option in options:
            result += option + '\n'
        result += '答案：' + answer + '\n\n'
    return result


def text_to_excel(txt):
    '''将txt转换为为Excel表格'''
    with open(txt, 'r', encoding='utf-8') as f:
        lines = f.readlines()
    sheets = {}
    sheet_name = ''
    for line in lines:
        if line.startswith('#'):
            sheet_name = line[1:].strip()
            sheets[sheet_name] = []
        else:
            sheets[sheet_name].append(line.strip())
    writer = pd.ExcelWriter('RainClass_PPT_Result.xlsx', engine='xlsxwriter')
    for sheet_name, rows in sheets.items():
        data = []
        for row in rows:
            if row.startswith('(单选)') or row.startswith('(多选)'):
                question = row
                options = []
                answer = ''
            elif row.startswith('A.') or row.startswith('B.') or row.startswith('C.') or row.startswith(
                    'D.') or row.startswith('E.'):
                options.append(row)
            elif row.startswith('答案：'):
                answer = row.split('：')[1]
                data.append([question] + options + [''] * (5 - len(options)) + ['', answer, ''])
        df = pd.DataFrame(data,
                          columns=['题目', '选项A', '选项B', '选项C', '选项D', '选项E', '您的答案', '正确答案', '正误'])
        df.to_excel(writer, sheet_name=sheet_name, index=False)
    writer.close()


def run():
    global exerciseNum
    if os.path.exists("RainClass_PPT_Result.xlsx"):
        del_file = input("\n需要删除上次生成结果文件RainClass_PPT_Result.xlsx才能继续\n删除吗？\n\tY/y：删除\n\t其他字符：退出\n请选择：").lower()
        if del_file == "y":
            os.remove("RainClass_PPT_Result.xlsx")
            print("已删除")
        else:
            print("您已取消练习题导出")
            return

    if not os.path.exists("./PPT"):
        os.mkdir("./PPT")
    choose = input('\n请将需要转换的PPT文件放在本程序同目录的PPT目录下\n\tY/y: 继续操作\n\t其他任意字符: 退出\n请选择：')

    if choose.lower() == 'y':
        if os.path.exists(".temp"):
            os.remove(".temp")

        print('\n开始读取./PPT目录中的文件\n')
        files = os.listdir('./PPT')
        for file in files:
            print('读取：\"' + file + '\"')
            text = read_pptx_text('./PPT/' + file, file.split('.')[0])
            with open('.temp', 'a', encoding='utf-8') as f1:
                f1.write(text)
                f1.close()
        print('所有PPT读取完成，共获取到' + str(exerciseNum) + '道题目\n')

        # 格式化获取到的文本，并导入到“RainClass_PPT_Result.xlsx”表格中
        with open('.temp', 'r', encoding='utf-8') as f2:
            data = f2.read()
            data = format_text(data)
            with open('.temp', 'w', encoding='utf-8') as f3:
                f3.write(data)
                f3.close()
        text_to_excel('./.temp')
        print('转换完成！结果已导出为RainClass_PPT_Result.xlsx\n返回主菜单')
        os.remove('./.temp')
        return
    else:
        print('您已取消练习题导出')
        return
